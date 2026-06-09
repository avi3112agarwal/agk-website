from pptx import Presentation

import sys
fn = sys.argv[1] if len(sys.argv) > 1 else "AGK_Client_Onboarding_v2.pptx"
p = Presentation(fn)
print(f"File: {fn}")
print(f"Total slides: {len(p.slides)}")
print("---")
for i, slide in enumerate(p.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip().replace("\n", " | ")
            if t:
                texts.append(t)
    headline = texts[0] if texts else "(empty)"
    subhead = texts[1] if len(texts) > 1 else ""
    print(f"Slide {i:>2}: {headline[:90]}")
    if subhead:
        print(f"         {subhead[:90]}")
print("---")
# Check for placeholder leftovers
import re
suspect = re.compile(r"\bx{3,}\b|lorem|ipsum|\bTODO\b|\[insert", re.IGNORECASE)
issues = []
for i, slide in enumerate(p.slides, 1):
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            for line in t.splitlines():
                if suspect.search(line):
                    issues.append(f"Slide {i}: {line!r}")
if issues:
    print("PLACEHOLDER LEFTOVERS:")
    for x in issues:
        print(" ", x)
else:
    print("OK — no placeholder leftovers found.")
